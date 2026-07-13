#!/usr/bin/env python3
"""카페 안도 × 인플루언서 TOP 5 — PPTX 빌더 (document-skills pptx 방식)"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 브랜드 컬러 (my_cafe.md) ─────────────────────────────
CREAM   = RGBColor(0xF5, 0xF0, 0xE8)
WALNUT  = RGBColor(0x5C, 0x46, 0x33)
SAGE    = RGBColor(0x8A, 0x9A, 0x7B)
SAGE_D  = RGBColor(0x6E, 0x7E, 0x60)
WARMGRAY= RGBColor(0xB8, 0xAC, 0x9E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x3A, 0x2E, 0x22)
KFONT   = "Apple SD Gothic Neo"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

def set_font(run, size=18, bold=False, color=INK, name=KFONT, italic=False):
    f = run.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, name
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def add_bg(slide, color=CREAM):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r

def add_rect(slide, x, y, w, h, color, rounded=True, line=None, radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line: shp.line.color.rgb = line; shp.line.width = Pt(1)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4):
    """lines: list of list-of-(text, kwargs) runs per paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        for text, kw in runs:
            r = p.add_run(); r.text = text; set_font(r, **kw)
    return tb

def header(slide, kicker, title):
    add_rect(slide, Inches(0.55), Inches(0.42), Inches(0.14), Inches(0.86), SAGE, rounded=False)
    add_text(slide, Inches(0.85), Inches(0.38), Inches(11.8), Inches(0.4),
             [[(kicker, dict(size=13, bold=True, color=SAGE_D))]])
    add_text(slide, Inches(0.85), Inches(0.68), Inches(11.8), Inches(0.7),
             [[(title, dict(size=30, bold=True, color=WALNUT))]])

# ── 데이터 ───────────────────────────────────────────────
TOP5 = [
    (1, "쑤현", "@ssu_h_k", "카페투어 · 디저트", 23004, "554", "21.1", "2.50%",
     "범위 내 최고 ER — 감성 카페·디저트 톤 일치, #쑤현_지역명 검색 태그로 유입이 계속 남는다"),
    (2, "재누", "@jaenu._", "서울 코스 · 공간", 1806, "599", "46.8", "35.76%",
     "'(저장) 코스 추천' 저장형 포맷 — 팔로워의 1/3이 좋아요. 최소 비용 와일드카드 ⚡"),
    (3, "픽철", "@picchul_", "카페·공간 사진가", 43478, "317", "7.6", "0.75%",
     "'골목 사랑방 · 조용한 공간' 서사 전문 — 2층 쇼룸 카페 스토리 최적, 사진 2차 활용 가치"),
    (4, "카페에디터", "@cafe__editor", "신상카페 소식", 52895, "497", "8.7", "0.96%",
     "신상·가오픈 카페 전문 지면 = 오픈 4주차인 우리에게 정확한 채널 (상한 소폭 초과 ⚠️)"),
    (5, "커푸얼", "@cupofpour", "카페 · 공간", 1003, "137", "26.9", "16.34%",
     "'한적한 혼자 시간' 감성이 브랜드 결 그 자체 — 무상 초대 씨딩 1순위 나노 ⚡"),
]
ER_CHART = [  # (이름, 팔로워표기, ER, top5여부)
    ("재누", "1.8k", 35.76, True), ("커푸얼", "1k", 16.34, True), ("쑤현", "23k", 2.50, True),
    ("지안", "6.2k", 1.14, False), ("카페에디터", "52.9k", 0.96, True), ("범스푼", "9.8k", 0.77, False),
    ("픽철", "43.5k", 0.75, True), ("방탱이", "151k", 0.40, False),
    ("성수교과서", "173k", 0.37, False), ("커피뚜벅이", "106k", 0.23, False),
]

# ═════════ Slide 1 · 표지 ═════════
s = prs.slides.add_slide(BLANK); add_bg(s)
add_rect(s, 0, Inches(6.9), SW, Inches(0.6), SAGE, rounded=False)
add_rect(s, Inches(0.9), Inches(1.55), Inches(0.18), Inches(2.5), SAGE, rounded=False)
add_text(s, Inches(1.35), Inches(1.45), Inches(11), Inches(0.5),
         [[("카페 안도 · 홍보 미션  |  [홍보] 우리 카페 띄워줄 인스타 인플루언서 찾기", dict(size=15, bold=True, color=SAGE_D))]])
add_text(s, Inches(1.35), Inches(1.95), Inches(11.2), Inches(2.2),
         [[("결이 맞는 인플루언서 TOP 5,", dict(size=44, bold=True, color=WALNUT))],
          [("인게이지먼트 실측으로 골랐습니다", dict(size=44, bold=True, color=WALNUT))]])
add_text(s, Inches(1.35), Inches(4.35), Inches(11), Inches(1.6),
         [[("2026-07-13 · Playwright(본인 로그인 세션) 실측 · 해시태그 6종 + 경쟁 카페 태그 2종 · 후보 12명 지표 분석", dict(size=14, color=INK))],
          [("☕ 카페 안도 — 성수 연무장길 2층 12평 · 자재 쇼룸 겸 카페 · @cafe.ando_ (팔로워 320)", dict(size=14, color=INK))],
          [("기준(my_cafe.md): 공간/카페/디저트 · 팔로워 5천~5만 마이크로", dict(size=14, bold=True, color=SAGE_D))]])

# ═════════ Slide 2 · 어떻게 찾았나 ═════════
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "METHOD", "어떻게 찾고, 어떻게 쟀나")
cards = [
    ("① 해시태그 발굴", "#성수카페 #성수동카페 #동네카페\n#분위기좋은카페 #카페추천 #디저트맛집\n→ 상위 게시물 ~140건에서 작성자 수집"),
    ("② 경쟁 카페의 손님들", "'메종 드 성수'·'로우키'는 실계정 없음\n→ 실존 성수 대형 카페 #옹근달 #킨포크성수\n게시물 작성자 = 성수 카페를 찍는 사람들"),
    ("③ 지표 실측", "후보 12명 프로필 + 최근 게시물 9개의\n좋아요·댓글 수집 (고정글·좋아요 숨김 제외)\n→ 브랜드/맞팔 계정 걸러냄"),
]
for i, (t, body) in enumerate(cards):
    x = Inches(0.7 + i * 4.1)
    add_rect(s, x, Inches(1.75), Inches(3.85), Inches(2.75), WHITE, line=WARMGRAY)
    add_text(s, x + Inches(0.25), Inches(1.98), Inches(3.4), Inches(0.5),
             [[(t, dict(size=17, bold=True, color=SAGE_D))]])
    add_text(s, x + Inches(0.25), Inches(2.52), Inches(3.4), Inches(1.9),
             [[(ln, dict(size=12.5, color=INK))] for ln in body.split("\n")], space_after=6)
add_rect(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.9), WALNUT)
add_text(s, Inches(1.1), Inches(5.1), Inches(11.2), Inches(0.6),
         [[("인게이지먼트율(ER) = (최근 9개 평균 좋아요 + 평균 댓글) ÷ 팔로워 × 100", dict(size=19, bold=True, color=CREAM))]])
add_text(s, Inches(1.1), Inches(5.78), Inches(11.2), Inches(0.9),
         [[("팔로워는 살 수 있어도 ER은 못 산다 — 팔로워 수보다 '내 게시물에 실제로 반응하는 사람'과 우리 타깃과의 결이 핵심", dict(size=13.5, color=CREAM))]])

# ═════════ Slide 3 · 핵심 발견 (ER 차트) ═════════
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "KEY FINDING", "팔로워 수 ≠ 효과 — 대형 0.2~0.4% vs 마이크로·나노 2.5~35%")
top_y, row_h = 1.85, 0.47
max_bar = 5.8
max_er = max(e for _, _, e, _ in ER_CHART)
for i, (name, fw, er, is5) in enumerate(ER_CHART):
    y = Inches(top_y + i * row_h)
    add_text(s, Inches(0.7), y, Inches(2.15), Inches(0.35),
             [[(f"{name} ({fw})", dict(size=13, bold=is5, color=WALNUT if is5 else WARMGRAY))]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    w = max(0.25, math.sqrt(er / max_er) * max_bar)
    add_rect(s, Inches(3.0), y + Inches(0.03), Inches(w), Inches(0.30),
             SAGE if is5 else WARMGRAY, radius=0.5)
    add_text(s, Inches(3.0 + w + 0.12), y, Inches(1.6), Inches(0.35),
             [[(f"{er:.2f}%", dict(size=13, bold=True, color=SAGE_D if is5 else WARMGRAY))]],
             anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(10.05), Inches(1.9), Inches(2.85), Inches(3.2), WHITE, line=WARMGRAY)
add_text(s, Inches(10.3), Inches(2.1), Inches(2.4), Inches(2.9),
         [[("읽는 법", dict(size=14, bold=True, color=SAGE_D))],
          [("■ 초록 = TOP 5", dict(size=12, color=SAGE_D))],
          [("■ 회색 = 비교군", dict(size=12, color=WARMGRAY))],
          [("", dict(size=6))],
          [("10만+ 대형 3계정 평균 ER 0.33% — 같은 예산이면 결 맞는 마이크로 여러 명이 8~100배 효율", dict(size=12, color=INK))]], space_after=6)
add_text(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.4),
         [[("* 막대 길이는 √스케일(수치 참조) · 범스푼(9.8k)·지안(6.2k)은 맞팔로 키운 팔로워라 ER 저조 → 탈락 · 2026-07-13 실측", dict(size=11, color=WARMGRAY))]])

# ═════════ Slide 4 · TOP 5 표 ═════════
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "SHORTLIST", "TOP 5 — 인게이지먼트율 + 적합도 종합")
rows, cols = 6, 6
tbl = s.shapes.add_table(rows, cols, Inches(0.55), Inches(1.8), Inches(12.25), Inches(4.5)).table
widths = [2.35, 1.75, 1.15, 1.35, 1.05, 4.60]
for c, wd in enumerate(widths): tbl.columns[c].width = Inches(wd)
heads = ["계정", "분야", "팔로워", "평균 ❤ / 💬", "ER", "카페 안도 적합도"]
for c, h in enumerate(heads):
    cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = WALNUT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h; set_font(r, size=13.5, bold=True, color=CREAM)
for i, (rk, nm, handle, cat, fw, al, ac, er, fit) in enumerate(TOP5):
    vals = [f"{rk}. {nm}  {handle}", cat, f"{fw:,}", f"{al} / {ac}", er, fit]
    for c, v in enumerate(vals):
        cell = tbl.cell(i + 1, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else CREAM
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left, cell.margin_right = Inches(0.08), Inches(0.06)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c in (0, 5) else PP_ALIGN.CENTER
        r = p.add_run(); r.text = v
        set_font(r, size=11.5 if c == 5 else 12.5,
                 bold=(c == 0 or c == 4), color=SAGE_D if c == 4 else INK)
tbl.rows[0].height = Inches(0.5)
for i in range(1, 6): tbl.rows[i].height = Inches(0.78)
add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
         [[("⚡ 재누·커푸얼은 5천 미만 나노지만 ER 압도적 → 와일드카드 선정 · ⚠️ 카페에디터는 상한 5만을 2.9k 초과(예외 추천) · 드림픽 성수교과서(173k)는 유료 대신 '로컬 신규 카페 제보' 무료 트랙", dict(size=11.5, color=WARMGRAY))]])

# ═════════ Slide 5~7 · TOP 3 상세 + DM ═════════
DETAILS = [
    ("TOP 1 · 본진", "쑤현  @ssu_h_k", "23,004 팔로워 · 평균 ❤ 554 · 💬 21.1 · ER 2.50%",
     ["카페·디저트 전문 감성 계정 — 사진 톤(따뜻한 우드·디저트 클로즈업)이 카페 안도 무드와 일치",
      "#쑤현_지역명 검색 태그 운영 → '#쑤현_성수' 게시물은 검색 유입이 계속 남는 자산형 노출",
      "제안: 흑임자 크림라떼 + 무화과 바스크 페어링 초대 (동행 1인 풀코스) + 원고료 협의"],
     "쑤현님 안녕하세요! 성수에서 작은 카페를 하는 '카페 안도' 사장입니다 ☕\n#쑤현_성수 태그로 성수 카페들 둘러보다가, 디저트 사진에 담기는 따뜻한 톤이\n저희 공간이랑 정말 잘 맞겠다 싶어 용기 내 연락드려요. (…전문은 influencers.md)"),
    ("TOP 2 · 가성비 와일드카드", "재누  @jaenu._", "1,806 팔로워 · 평균 ❤ 599 · 💬 46.8 · ER 35.76%",
     ["'(저장) 코스 추천' 저장형 포맷 — 탐색탭 도달이 팔로워를 수십 배 초과 (좋아요가 팔로워의 1/3)",
      "성수 대형 베이커리 카페 게시 이력 → 반대 컨셉 '성수 조용한 2층 숨은 카페 코스' 소재로 제안",
      "제안: 동행 포함 전 메뉴 무료 + 소정의 제작비 — 나노라 비용 최소, 저장형이라 수명 김"],
     "재누님 안녕하세요, 성수 골목 2층에 있는 '카페 안도'입니다 :)\n\"(저장) 코스 추천\" 시리즈 너무 잘 보고 있어요. 성수 대형 베이커리 카페 편 보고,\n저희 같은 반대편 컨셉(조용한 소형)도 코스 소재가 되겠다 싶어 연락드려요. (…)"),
    ("TOP 3 · 브랜드 서사", "픽철  @picchul_", "43,478 팔로워 · 평균 ❤ 317 · 💬 7.6 · ER 0.75%",
     ["'골목 사랑방 · 13년째 한자리 로스터리' 같은 공간 서사 전문 사진가 — 우리 스토리에 정확히 꽂히는 시선",
      "'인테리어 디자이너가 도면부터 직접 지은 살아있는 포트폴리오 카페' = 픽철 포맷 그 자체",
      "제안: 조용한 시간대 촬영 초대 + 원고료 + 촬영 사진의 우리 계정 2차 활용 별도 협의 (사진 자산 확보)"],
     "픽철님 안녕하세요. 성수 연무장길 안쪽 골목에서 '카페 안도'라는 2층 카페를 하는\n사장이자 7년째 인테리어 회사를 운영하는 디자이너입니다. 공간을 서사로 읽어주시는\n시선이, 저희 카페가 꼭 한번 담기고 싶은 시선이라 연락드렸어요. (…)"),
]
for kicker, name, stats, whys, dm in DETAILS:
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, kicker, name)
    add_rect(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.62), SAGE)
    add_text(s, Inches(1.0), Inches(1.78), Inches(11.4), Inches(0.45),
             [[(stats, dict(size=16, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.75), Inches(2.6), Inches(6.3), Inches(0.4),
             [[("왜 이 사람인가", dict(size=15, bold=True, color=SAGE_D))]])
    add_text(s, Inches(0.75), Inches(3.05), Inches(6.35), Inches(3.6),
             [[("·  " + w, dict(size=13, color=INK))] for w in whys], space_after=10)
    add_rect(s, Inches(7.35), Inches(2.6), Inches(5.35), Inches(4.15), WHITE, line=SAGE)
    add_text(s, Inches(7.6), Inches(2.8), Inches(4.9), Inches(0.4),
             [[("💌 컨택 DM 초안 (발췌)", dict(size=14, bold=True, color=SAGE_D))]])
    add_text(s, Inches(7.6), Inches(3.25), Inches(4.9), Inches(3.3),
             [[(ln, dict(size=11.5, color=INK))] for ln in dm.split("\n")], space_after=7)

# ═════════ Slide 8 · 실행 플랜 ═════════
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "NEXT", "DM 보내기 전 실행 플랜")
steps = [
    ("1", "우리 계정 먼저 정비", "@cafe.ando_ 게시물 9개·릴스 0개로 DM 보내면 회신율 급락 — 샘플월·시그니처 릴스 3개 먼저"),
    ("2", "하루 1~2명씩 발송", "동시 다발 DM = 스팸 인상 + 계정 제재 리스크. 쑤현 → 재누 → 픽철 순"),
    ("3", "1주 무응답 시", "스토리 멘션으로 한 번만 소프트 리마인드 (재촉 금지)"),
    ("4", "무료 트랙 병행", "성수교과서(173k)에 협찬 아닌 '로컬 신규 카페 제보' DM — Supporting Local Business 표방 계정"),
    ("5", "방문일 준비", "주말 한정 무화과 바스크 예약 확보 + 촬영 시간대 조용히 비워두기"),
]
for i, (n, t, d) in enumerate(steps):
    y = Inches(1.8 + i * 1.02)
    add_rect(s, Inches(0.7), y, Inches(0.62), Inches(0.62), SAGE, radius=0.5)
    add_text(s, Inches(0.7), y, Inches(0.62), Inches(0.62),
             [[(n, dict(size=20, bold=True, color=WHITE))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), y - Inches(0.02), Inches(11.2), Inches(0.4),
             [[(t, dict(size=16, bold=True, color=WALNUT))]])
    add_text(s, Inches(1.55), y + Inches(0.36), Inches(11.2), Inches(0.4),
             [[(d, dict(size=12.5, color=INK))]])
add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
         [[("전체 데이터: week-6/influencer/influencers.md · candidates.json  |  목표(my_cafe.md §9): 인스타 320 → 3,000 팔로워", dict(size=11.5, color=WARMGRAY))]])

OUT = "/Users/pyounghwahong/AFM/week-6/influencer/influencers.pptx"
prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
